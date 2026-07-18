"""
Local, dependency-light document parser.

Handles PDF, DOCX, PPTX, XLSX, Markdown and HTML without any external
service call. This is the default parser and the fallback whenever
LlamaParse is not configured (no API key) or fails.

Layout fidelity is intentionally simpler than LlamaParse: headings are
inferred from font size / style heuristics where the library exposes them,
otherwise every page becomes one "text" element plus separately extracted
tables where the underlying library supports it.
"""

from __future__ import annotations

from pathlib import Path

from app.parsing.interfaces import DocumentParser, ParsedDocument, ParsedElement


class LocalDocumentParser(DocumentParser):
    supported_extensions = {".pdf", ".docx", ".pptx", ".xlsx", ".md", ".html", ".htm", ".txt"}

    async def parse(self, file_path: str, document_id: str) -> ParsedDocument:
        path = Path(file_path)
        suffix = path.suffix.lower()

        if suffix == ".pdf":
            elements = self._parse_pdf(path)
        elif suffix == ".docx":
            elements = self._parse_docx(path)
        elif suffix == ".pptx":
            elements = self._parse_pptx(path)
        elif suffix == ".xlsx":
            elements = self._parse_xlsx(path)
        elif suffix in {".md", ".txt"}:
            elements = self._parse_markdown(path)
        elif suffix in {".html", ".htm"}:
            elements = self._parse_html(path)
        else:
            raise ValueError(f"Unsupported extension: {suffix}")

        return ParsedDocument(document_id=document_id, filename=path.name, elements=elements)

    def _parse_pdf(self, path: Path) -> list[ParsedElement]:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        elements: list[ParsedElement] = []
        for page_number, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                elements.append(ParsedElement(kind="text", content=text.strip(), page=page_number))
        return elements

    def _parse_docx(self, path: Path) -> list[ParsedElement]:
        import docx

        document = docx.Document(str(path))
        elements: list[ParsedElement] = []
        for paragraph in document.paragraphs:
            if not paragraph.text.strip():
                continue
            style = (paragraph.style.name or "").lower()
            kind = "heading" if "heading" in style else "text"
            elements.append(ParsedElement(kind=kind, content=paragraph.text.strip()))
        for table_index, table in enumerate(document.tables):
            rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            rendered = "\n".join(" | ".join(row) for row in rows)
            elements.append(
                ParsedElement(kind="table", content=rendered, metadata={"table_index": table_index})
            )
        return elements

    def _parse_pptx(self, path: Path) -> list[ParsedElement]:
        from pptx import Presentation

        presentation = Presentation(str(path))
        elements: list[ParsedElement] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    texts.append(shape.text_frame.text.strip())
            if texts:
                elements.append(
                    ParsedElement(kind="text", content="\n".join(texts), page=slide_number)
                )
        return elements

    def _parse_xlsx(self, path: Path) -> list[ParsedElement]:
        from openpyxl import load_workbook

        workbook = load_workbook(str(path), data_only=True)
        elements: list[ParsedElement] = []
        for sheet in workbook.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    rows.append(" | ".join("" if cell is None else str(cell) for cell in row))
            if rows:
                elements.append(
                    ParsedElement(kind="table", content="\n".join(rows), section=sheet.title)
                )
        return elements

    def _parse_markdown(self, path: Path) -> list[ParsedElement]:
        text = path.read_text(encoding="utf-8", errors="ignore")
        elements: list[ParsedElement] = []
        current_section = None
        buffer: list[str] = []

        def flush():
            if buffer:
                elements.append(
                    ParsedElement(kind="text", content="\n".join(buffer).strip(), section=current_section)
                )
                buffer.clear()

        for line in text.splitlines():
            if line.startswith("#"):
                flush()
                current_section = line.lstrip("#").strip()
                elements.append(ParsedElement(kind="heading", content=current_section, section=current_section))
            else:
                buffer.append(line)
        flush()
        return elements

    def _parse_html(self, path: Path) -> list[ParsedElement]:
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="ignore"), "html.parser")
        elements: list[ParsedElement] = []
        for tag in soup.find_all(["h1", "h2", "h3", "p", "li", "table"]):
            text = tag.get_text(strip=True, separator=" ")
            if not text:
                continue
            if tag.name in {"h1", "h2", "h3"}:
                elements.append(ParsedElement(kind="heading", content=text))
            elif tag.name == "table":
                elements.append(ParsedElement(kind="table", content=text))
            else:
                elements.append(ParsedElement(kind="text", content=text))
        return elements
